from pathlib import Path
import json
from datetime import datetime

from PIL import Image

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = Path(__file__).resolve().parents[1]

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT
}

def load_json(path):
    with open(path, "r", encoding="utf-8-sig") as f:
        return json.load(f)

def resolve_style(el, theme_styles):
    style = {}

    style_ref = el.get("style_ref")
    if style_ref and style_ref in theme_styles:
        style.update(theme_styles[style_ref])

    if "style" in el:
        style.update(el["style"])

    return style

def apply_text_style(paragraph, run, style):
    font_family = style.get("font_family", "Yu Gothic")
    font_size = style.get("font_size", 18)
    bold = style.get("bold", False)
    italic = style.get("italic", False)
    color = style.get("color", "333333")
    align = style.get("align", "left")

    run.font.name = font_family
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic

    try:
        run.font.color.rgb = RGBColor.from_string(color.replace("#", ""))
    except Exception:
        run.font.color.rgb = RGBColor(51, 51, 51)

    paragraph.alignment = ALIGN_MAP.get(align, PP_ALIGN.LEFT)

def add_textbox(slide, el, theme_styles):
    style = resolve_style(el, theme_styles)

    box = slide.shapes.add_textbox(
        Inches(el["x"]),
        Inches(el["y"]),
        Inches(el["w"]),
        Inches(el["h"])
    )

    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    # Keep PPTX text drawing aligned with detected bbox.
    # PowerPoint text boxes have internal margins; for OCR-derived boxes,
    # even 0.02 inch can visibly shift small Japanese text.
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    lines = el.get("text", "").split("\\n")

    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        apply_text_style(p, run, style)

def create_faded_background(src_path, slide_id, opacity):
    """
    opacity:
      0.0 = 真っ白
      1.0 = 元画像そのまま
    """
    cache_dir = BASE_DIR / "output" / "_cache"
    cache_dir.mkdir(parents=True, exist_ok=True)

    img = Image.open(src_path).convert("RGB")
    white = Image.new("RGB", img.size, "white")

    opacity = max(0.0, min(1.0, float(opacity)))
    faded = Image.blend(white, img, opacity)

    out_path = cache_dir / f"{slide_id}_faded_bg_{int(opacity * 100):03d}.png"
    faded.save(out_path)

    return out_path

def add_image(slide, el, img_path_override=None):
    img_path = img_path_override if img_path_override else BASE_DIR / el["path"]

    if not img_path.exists():
        print(f"[WARN] image not found: {img_path}")
        return

    slide.shapes.add_picture(
        str(img_path),
        Inches(el["x"]),
        Inches(el["y"]),
        width=Inches(el["w"]),
        height=Inches(el["h"])
    )



def load_assets_manifest(slide_id):
    path = BASE_DIR / "assets" / slide_id / "assets_manifest.json"
    if not path.exists():
        return {"slide_id": slide_id, "assets": []}

    try:
        data = load_json(path)
    except Exception as exc:
        print(f"[WARN] failed to load assets manifest: {path} ({exc})")
        return {"slide_id": slide_id, "assets": []}

    if not isinstance(data, dict):
        return {"slide_id": slide_id, "assets": []}

    assets = data.get("assets", [])
    if not isinstance(assets, list):
        data["assets"] = []

    return data


def find_reference_image_path(slide_spec):
    for el in slide_spec.get("elements", []):
        if isinstance(el, dict) and el.get("kind") == "image" and el.get("name") == "reference_bg":
            path = el.get("path")
            if path:
                return BASE_DIR / path
    return None


def bbox_px_to_inches(bbox_px, image_size, slide_spec):
    if not isinstance(bbox_px, list) or len(bbox_px) != 4:
        return None

    img_w, img_h = image_size
    if img_w <= 0 or img_h <= 0:
        return None

    canvas = slide_spec.get("canvas", {})
    slide_w = float(canvas.get("width_in", 13.333))
    slide_h = float(canvas.get("height_in", 7.5))

    try:
        x, y, w, h = [float(v) for v in bbox_px]
    except Exception:
        return None

    return {
        "x": (x / img_w) * slide_w,
        "y": (y / img_h) * slide_h,
        "w": (w / img_w) * slide_w,
        "h": (h / img_h) * slide_h,
    }


def add_asset_picture(slide, asset, slide_spec):
    slide_id = slide_spec.get("slide_id")
    filename = asset.get("filename")
    bbox_px = asset.get("bbox_px")

    if not slide_id or not filename:
        return

    asset_path = BASE_DIR / "assets" / slide_id / filename
    if not asset_path.exists():
        print(f"[WARN] asset image not found: {asset_path}")
        return

    ref_path = find_reference_image_path(slide_spec)
    if not ref_path or not ref_path.exists():
        print(f"[WARN] reference image not found for asset placement: {slide_id}")
        return

    try:
        with Image.open(ref_path) as img:
            image_size = img.size
    except Exception as exc:
        print(f"[WARN] failed to read reference image size: {ref_path} ({exc})")
        return

    box = bbox_px_to_inches(bbox_px, image_size, slide_spec)
    if not box:
        print(f"[WARN] invalid asset bbox: {slide_id} {filename}")
        return

    slide.shapes.add_picture(
        str(asset_path),
        Inches(box["x"]),
        Inches(box["y"]),
        width=Inches(box["w"]),
        height=Inches(box["h"])
    )

    print(f"asset: {slide_id} {filename}")


def add_pptx_assets(slide, slide_spec):
    slide_id = slide_spec.get("slide_id")
    if not slide_id:
        return

    manifest = load_assets_manifest(slide_id)
    assets = manifest.get("assets", [])

    for asset in assets:
        if not isinstance(asset, dict):
            continue
        if not bool(asset.get("use_in_pptx", False)):
            continue
        add_asset_picture(slide, asset, slide_spec)



def add_shape(slide, el):
    shape_type = el.get("shape_type", "rect")

    if shape_type != "rect":
        print(f"[WARN] unsupported shape_type: {shape_type}")
        return

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(el["x"]),
        Inches(el["y"]),
        Inches(el["w"]),
        Inches(el["h"])
    )

    fill = el.get("fill")
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill.replace("#", ""))
    else:
        shape.fill.background()

    line = el.get("line")
    if line:
        shape.line.color.rgb = RGBColor.from_string(line.replace("#", ""))
    else:
        shape.line.fill.background()


def text_block_to_element(block):
    placement = block.get("placement_in", {})

    if not isinstance(placement, dict):
        placement = {}

    x = placement.get("x", block.get("x", 0.5))
    y = placement.get("y", block.get("y", 0.5))
    w = placement.get("w", placement.get("width", block.get("w", block.get("width", 4.0))))
    h = placement.get("h", placement.get("height", block.get("h", block.get("height", 0.4))))

    return {
        "kind": "text",
        "name": block.get("id", "ocr_text"),
        "text": block.get("text", ""),
        "x": float(x),
        "y": float(y),
        "w": float(w),
        "h": float(h),
        "style_ref": block.get("style_ref", block.get("semantic_role", "p")),
        "semantic_role": block.get("semantic_role", "p"),
        "source": block.get("source", "text_blocks"),
        "text_mode": block.get("text_mode", "ocr")
    }


def resolve_slide_elements(slide_spec):
    elements = slide_spec.get("elements", [])

    if not isinstance(elements, list):
        elements = []

    text_blocks = slide_spec.get("text_blocks", [])

    if not isinstance(text_blocks, list) or len(text_blocks) == 0:
        return elements

    non_text_elements = [
        el for el in elements
        if isinstance(el, dict) and el.get("kind") != "text"
    ]

    ocr_text_elements = [
        text_block_to_element(block)
        for block in text_blocks
        if isinstance(block, dict) and str(block.get("text", "")).strip()
    ]

    return non_text_elements + ocr_text_elements



def normalize_hex_color(value, fallback="333333"):
    if not value:
        return fallback
    color = str(value).replace("#", "").strip()
    if len(color) == 6:
        return color.upper()
    return fallback


def find_style_color(theme_styles, style_keys, fallback="333333"):
    for key in style_keys:
        style = theme_styles.get(key, {})
        if isinstance(style, dict) and style.get("color"):
            return normalize_hex_color(style.get("color"), fallback)
    return fallback


def make_theme_tokens(theme_styles):
    """
    Build practical theme tokens from existing role styles.
    This keeps the Pattern Sheet independent from future token schemas.
    """
    main = find_style_color(theme_styles, ["left.h1", "main.h1", "title.h1"], "243B53")
    sub = find_style_color(theme_styles, ["left.h2", "main.h2", "title.h2"], "486581")
    body = find_style_color(theme_styles, ["left.p", "main.p", "body.p"], "334E68")
    note = find_style_color(theme_styles, ["footer.note", "meta.small"], "627D98")

    return {
        "main": main,
        "sub": sub,
        "accent": "2F80ED",
        "background_light": "FFFFFF",
        "background_dark": "102A43",
        "surface_light": "F0F4F8",
        "surface_dark": "243B53",
        "text_on_light": body,
        "text_on_dark": "FFFFFF",
        "muted_text_on_light": note,
        "muted_text_on_dark": "D9E2EC",
        "border_on_light": "BCCCDC",
        "border_on_dark": "829AB1",
        "highlight": "F2C94C",
        "warning": "F2994A",
        "success": "27AE60",
        "danger": "EB5757",
    }


def rgb(hex_color):
    return RGBColor.from_string(normalize_hex_color(hex_color))



def disable_shape_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass
    return shape

def add_pattern_label(slide, text, x, y, w, h=0.18, color="627D98", size=7, bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Yu Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    p.alignment = PP_ALIGN.LEFT
    return box


def add_pattern_box(
    slide,
    text,
    x,
    y,
    w,
    h,
    fill,
    line,
    font,
    font_size=8,
    bold=False,
    radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE,
):
    shape = slide.shapes.add_shape(
        radius_shape,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    disable_shape_shadow(shape)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)

    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Yu Gothic"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(font)
    p.alignment = PP_ALIGN.CENTER

    return shape


def add_pattern_table(slide, title, x, y, w, h, tokens, pattern):
    add_pattern_label(slide, title, x, y - 0.18, w, bold=True)

    rows = 4
    cols = 4
    cell_w = w / cols
    cell_h = h / rows

    for r in range(rows):
        for c in range(cols):
            cx = x + c * cell_w
            cy = y + r * cell_h

            is_header = r == 0
            is_alt = r % 2 == 0
            is_highlight = pattern == "highlight" and r == 2 and c == 2

            fill = tokens["background_light"]
            line = tokens["border_on_light"]
            font = tokens["text_on_light"]
            line_width = 0.4

            if pattern == "header" and is_header:
                fill = tokens["main"]
                font = tokens["text_on_dark"]
                line = tokens["main"]

            elif pattern == "simple" and is_header:
                fill = tokens["surface_light"]

            elif pattern == "horizontal":
                line = tokens["border_on_light"]

            elif pattern == "zebra" and is_alt and not is_header:
                fill = tokens["surface_light"]

            elif pattern == "highlight" and is_highlight:
                fill = tokens["accent"]
                font = tokens["text_on_dark"]
                line = tokens["accent"]

            elif pattern == "minimal":
                line = tokens["background_light"]
                line_width = 0.1
                if is_header:
                    font = tokens["main"]

            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(cx),
                Inches(cy),
                Inches(cell_w),
                Inches(cell_h),
            )
            disable_shape_shadow(cell)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill)
            cell.line.color.rgb = rgb(line)
            cell.line.width = Pt(line_width)

            if pattern == "horizontal":
                cell.line.color.rgb = rgb(tokens["background_light"])
                cell.line.width = Pt(0.1)
                if r > 0:
                    rule = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(cx),
                        Inches(cy),
                        Inches(cell_w),
                        Inches(0.006),
                    )
                    disable_shape_shadow(rule)
                    rule.fill.solid()
                    rule.fill.fore_color.rgb = rgb(tokens["border_on_light"])
                    rule.line.color.rgb = rgb(tokens["border_on_light"])

            tf = cell.text_frame
            tf.clear()
            tf.margin_left = Inches(0.02)
            tf.margin_right = Inches(0.02)
            tf.margin_top = Inches(0.01)
            tf.margin_bottom = Inches(0.01)

            p = tf.paragraphs[0]
            run = p.add_run()
            if is_header:
                run.text = f"H{c+1}"
            else:
                run.text = f"{r}.{c+1}"
            run.font.name = "Yu Gothic"
            run.font.size = Pt(5.5)
            run.font.bold = is_header
            run.font.color.rgb = rgb(font)
            p.alignment = PP_ALIGN.CENTER



def add_arrow_icon(slide, text, shape_type, x, y, w, h, fill, line, font, font_size=6, bold=True):
    shape = slide.shapes.add_shape(
        shape_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    disable_shape_shadow(shape)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.75)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Yu Gothic"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(font)
    p.alignment = PP_ALIGN.CENTER

    return shape









def set_table_cell_text(cell, text, fill, font, font_size=5.5, bold=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(fill)

    tf = cell.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Yu Gothic"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(font)
    p.alignment = PP_ALIGN.CENTER


def add_reusable_pattern_table(slide, title, x, y, w, h, tokens, pattern):
    """
    Add a real PowerPoint table, not rectangle-cell mockups.
    The generated table can be selected and reused as a table in PowerPoint.
    """
    add_pattern_label(slide, title, x, y - 0.18, w, bold=True)

    rows = 4
    cols = 4
    table_shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    disable_shape_shadow(table_shape)

    table = table_shape.table

    for c in range(cols):
        table.columns[c].width = Inches(w / cols)

    for r in range(rows):
        table.rows[r].height = Inches(h / rows)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)

            is_header = r == 0
            is_alt = r % 2 == 0
            is_highlight = pattern == "highlight" and r == 2 and c == 2

            fill = tokens["background_light"]
            font = tokens["text_on_light"]
            bold = is_header

            if pattern == "simple":
                if is_header:
                    fill = tokens["surface_light"]
                    font = tokens["main"]

            elif pattern == "header":
                if is_header:
                    fill = tokens["main"]
                    font = tokens["text_on_dark"]
                else:
                    fill = tokens["background_light"]

            elif pattern == "horizontal":
                if is_header:
                    fill = tokens["background_light"]
                    font = tokens["main"]
                else:
                    fill = tokens["background_light"]

            elif pattern == "zebra":
                if is_header:
                    fill = tokens["main"]
                    font = tokens["text_on_dark"]
                elif is_alt:
                    fill = tokens["surface_light"]

            elif pattern == "highlight":
                if is_header:
                    fill = tokens["surface_light"]
                    font = tokens["main"]
                elif is_highlight:
                    fill = tokens["accent"]
                    font = tokens["text_on_dark"]
                    bold = True

            elif pattern == "minimal":
                fill = tokens["background_light"]
                if is_header:
                    font = tokens["main"]
                else:
                    font = tokens["muted_text_on_light"]

            if is_header:
                label = ["Item", "Plan", "Actual", "Status"][c]
            else:
                label = [
                    ["A", "100", "96", "OK"],
                    ["B", "120", "132", "Watch"],
                    ["C", "80", "74", "Low"],
                ][r - 1][c]

            set_table_cell_text(cell, label, fill, font, font_size=5.2, bold=bold)

    return table_shape





def add_process_flow_sample(slide, x, y, tokens):
    """
    Practical process flow sample for business slides.
    Simple two-lane layout without overlap.
    """
    add_pattern_label(slide, "Process Flow Sample", x, y - 0.20, 3.2, bold=True)

    step_w = 0.62
    step_h = 0.28
    arrow_w = 0.20
    gap = 0.10

    steps = [
        ("Input", tokens["surface_light"], tokens["border_on_light"], tokens["text_on_light"]),
        ("Review", tokens["main"], tokens["main"], tokens["text_on_dark"]),
        ("Approve", tokens["accent"], tokens["accent"], tokens["text_on_dark"]),
        ("Output", tokens["success"], tokens["success"], tokens["text_on_dark"]),
    ]

    cursor = x
    top_y = y
    for i, (label, fill, line, font) in enumerate(steps):
        add_pattern_box(
            slide,
            label,
            cursor,
            top_y,
            step_w,
            step_h,
            fill,
            line,
            font,
            5.3,
            True,
        )

        if i < len(steps) - 1:
            add_arrow_icon(
                slide,
                "",
                MSO_SHAPE.RIGHT_ARROW,
                cursor + step_w + 0.02,
                top_y + 0.06,
                arrow_w,
                0.16,
                tokens["border_on_light"],
                tokens["border_on_light"],
                tokens["text_on_dark"],
                1,
            )
        cursor += step_w + arrow_w + gap

    # lower lane: Review -> Decision -> Revise
    decision_y = y + 0.48
    decision_x = x + 0.92
    revise_x = decision_x + 0.92

    # vertical arrow from Review to Decision
    add_arrow_icon(
        slide,
        "",
        MSO_SHAPE.DOWN_ARROW,
        x + 1.04,
        y + 0.30,
        0.18,
        0.18,
        tokens["warning"],
        tokens["warning"],
        tokens["text_on_dark"],
        1,
    )

    add_pattern_box(
        slide,
        "Decision",
        decision_x,
        decision_y,
        0.72,
        0.26,
        tokens["warning"],
        tokens["warning"],
        tokens["text_on_dark"],
        5.0,
        True,
    )

    add_arrow_icon(
        slide,
        "",
        MSO_SHAPE.RIGHT_ARROW,
        decision_x + 0.74,
        decision_y + 0.05,
        0.18,
        0.16,
        tokens["danger"],
        tokens["danger"],
        tokens["text_on_dark"],
        1,
    )

    add_pattern_box(
        slide,
        "Revise",
        revise_x,
        decision_y,
        0.72,
        0.26,
        tokens["background_light"],
        tokens["danger"],
        tokens["danger"],
        5.0,
        True,
    )


def get_theme_style(theme_styles, style_keys, fallback=None):
    fallback = fallback or {}
    for key in style_keys:
        style = theme_styles.get(key, {})
        if isinstance(style, dict):
            merged = dict(fallback)
            merged.update(style)
            return merged
    return dict(fallback)



def resolve_font_size(style, fallback=10):
    """
    Resolve font size from theme style.
    Supports multiple key names for forward/backward compatibility.
    """
    for key in ("font_size", "fontSize", "font_size_pt", "fontSizePt", "size"):
        value = style.get(key)
        if value is None:
            continue
        try:
            return float(value)
        except Exception:
            continue
    return float(fallback)


def clamp_font_size_for_sample(font_size):
    """
    Keep the sample usable on a single appendix slide,
    while still reflecting the theme size differences.
    """
    try:
        size = float(font_size)
    except Exception:
        size = 10.0
    return max(6.0, min(size, 24.0))


def add_font_style_sample(slide, label, sample, x, y, w, h, style, tokens):
    """
    Add an actual-size font swatch.
    The sample text is bottom-aligned so larger fonts grow upward.
    """
    font_family = style.get("font_family", "Yu Gothic")
    font_size = resolve_font_size(style, 10)
    color = normalize_hex_color(style.get("color", tokens["text_on_light"]), tokens["text_on_light"])
    bold = bool(style.get("bold", False))
    italic = bool(style.get("italic", False))

    # Role label
    label_box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.12),
    )
    label_tf = label_box.text_frame
    label_tf.clear()
    label_tf.margin_left = Inches(0)
    label_tf.margin_right = Inches(0)
    label_tf.margin_top = Inches(0)
    label_tf.margin_bottom = Inches(0)

    label_p = label_tf.paragraphs[0]
    label_r = label_p.add_run()
    label_r.text = label
    label_r.font.name = "Yu Gothic"
    label_r.font.size = Pt(5.0)
    label_r.font.bold = True
    label_r.font.color.rgb = rgb(tokens["muted_text_on_light"])
    label_p.alignment = PP_ALIGN.CENTER

    # Actual-size sample, bottom-aligned.
    sample_box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y + 0.10),
        Inches(w),
        Inches(h - 0.26),
    )
    tf = sample_box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    try:
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    except Exception:
        pass

    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = sample
    r.font.name = font_family
    r.font.size = Pt(font_size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb(color)
    p.alignment = PP_ALIGN.CENTER

    # Font setting metadata under the baseline area.
    meta_box = slide.shapes.add_textbox(
        Inches(x),
        Inches(y + h - 0.14),
        Inches(w),
        Inches(0.14),
    )
    meta_tf = meta_box.text_frame
    meta_tf.clear()
    meta_tf.word_wrap = False
    meta_tf.margin_left = Inches(0)
    meta_tf.margin_right = Inches(0)
    meta_tf.margin_top = Inches(0)
    meta_tf.margin_bottom = Inches(0)

    meta_p = meta_tf.paragraphs[0]
    meta_r = meta_p.add_run()
    meta_r.text = f"{font_family} / {font_size:g}pt"
    meta_r.font.name = "Yu Gothic"
    meta_r.font.size = Pt(4.2)
    meta_r.font.color.rgb = rgb(tokens["muted_text_on_light"])
    meta_p.alignment = PP_ALIGN.CENTER

    return sample_box



def add_font_style_samples(slide, x, y, tokens, theme_styles):
    add_pattern_label(slide, "Font Styles", x, y - 0.18, 4.8, bold=True)

    fallback = {
        "font_family": "Yu Gothic",
        "font_size": 10,
        "color": tokens["text_on_light"],
        "bold": False,
        "italic": False,
    }

    samples = [
        ("H1", get_theme_style(theme_styles, ["left.h1", "main.h1", "title.h1"], {**fallback, "font_size": 13, "bold": True, "color": tokens["main"]})),
        ("H2", get_theme_style(theme_styles, ["left.h2", "main.h2", "title.h2"], {**fallback, "font_size": 11, "bold": True, "color": tokens["sub"]})),
        ("Body", get_theme_style(theme_styles, ["left.p", "main.p", "body.p"], {**fallback, "font_size": 9, "color": tokens["text_on_light"]})),
        ("Note", get_theme_style(theme_styles, ["footer.note", "meta.small"], {**fallback, "font_size": 8, "color": tokens["muted_text_on_light"]})),
    ]

    sample_text = "Aaあ1"
    box_w = 1.48
    gap = 0.18

    for i, (label, style) in enumerate(samples):
        add_font_style_sample(
            slide,
            label,
            sample_text,
            x + i * (box_w + gap),
            y,
            box_w,
            0.88,
            style,
            tokens,
        )


def add_theme_pattern_sheet(prs, theme_styles):
    tokens = make_theme_tokens(theme_styles)

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(tokens["background_light"])
    bg.line.color.rgb = rgb(tokens["background_light"])

    add_pattern_box(
        slide,
        "Theme Pattern Sheet",
        0.35,
        0.25,
        4.4,
        0.38,
        tokens["main"],
        tokens["main"],
        tokens["text_on_dark"],
        font_size=14,
        bold=True,
        radius_shape=MSO_SHAPE.RECTANGLE,
    )
    add_pattern_label(
        slide,
        "Reusable PowerPoint components generated from current theme styles",
        4.95,
        0.34,
        5.7,
        color=tokens["muted_text_on_light"],
        size=8,
    )

    add_pattern_label(slide, "Title Banner", 0.45, 0.9, 2.7, bold=True)
    add_pattern_box(slide, "Main Title", 0.45, 1.12, 2.7, 0.36, tokens["main"], tokens["main"], tokens["text_on_dark"], 9, True, MSO_SHAPE.RECTANGLE)
    add_pattern_box(slide, "Sub Title", 0.45, 1.56, 2.7, 0.36, tokens["sub"], tokens["sub"], tokens["text_on_dark"], 9, True, MSO_SHAPE.RECTANGLE)
    add_pattern_box(slide, "Light + Main", 0.45, 2.0, 2.7, 0.36, tokens["surface_light"], tokens["border_on_light"], tokens["main"], 9, True, MSO_SHAPE.RECTANGLE)

    add_pattern_label(slide, "Text Panel", 3.45, 0.9, 2.7, bold=True)
    add_pattern_box(slide, "Light panel text", 3.45, 1.12, 2.7, 0.58, tokens["background_light"], tokens["border_on_light"], tokens["text_on_light"], 8)
    add_pattern_box(slide, "Dark panel text", 3.45, 1.82, 2.7, 0.58, tokens["background_dark"], tokens["border_on_dark"], tokens["text_on_dark"], 8)
    add_pattern_box(slide, "Accent border", 3.45, 2.52, 2.7, 0.58, tokens["background_light"], tokens["accent"], tokens["text_on_light"], 8)

    add_pattern_label(slide, "Callout / Button", 6.45, 0.9, 2.7, bold=True)
    add_pattern_box(slide, "Accent", 6.45, 1.12, 1.2, 0.36, tokens["accent"], tokens["accent"], tokens["text_on_dark"], 8, True)
    add_pattern_box(slide, "Main", 7.85, 1.12, 1.2, 0.36, tokens["main"], tokens["main"], tokens["text_on_dark"], 8, True)
    add_pattern_box(slide, "Outline", 6.45, 1.64, 1.2, 0.36, tokens["background_light"], tokens["accent"], tokens["accent"], 8, True)
    add_pattern_box(slide, "Warning", 7.85, 1.64, 1.2, 0.36, tokens["warning"], tokens["warning"], tokens["text_on_dark"], 8, True)
    add_pattern_box(slide, "Success", 6.45, 2.16, 1.2, 0.36, tokens["success"], tokens["success"], tokens["text_on_dark"], 8, True)
    add_pattern_box(slide, "Danger", 7.85, 2.16, 1.2, 0.36, tokens["danger"], tokens["danger"], tokens["text_on_dark"], 8, True)

    add_pattern_label(slide, "Box Pattern", 9.65, 0.9, 2.9, bold=True)
    add_pattern_box(slide, "Border", 9.65, 1.12, 1.25, 0.52, tokens["background_light"], tokens["border_on_light"], tokens["text_on_light"], 7)
    add_pattern_box(slide, "Filled", 11.1, 1.12, 1.25, 0.52, tokens["surface_light"], tokens["surface_light"], tokens["text_on_light"], 7)
    add_pattern_box(slide, "Dark", 9.65, 1.82, 1.25, 0.52, tokens["surface_dark"], tokens["surface_dark"], tokens["text_on_dark"], 7)
    add_pattern_box(slide, "Accent", 11.1, 1.82, 1.25, 0.52, tokens["background_light"], tokens["accent"], tokens["accent"], 7)

    add_pattern_label(slide, "Chart Sample", 0.45, 3.35, 3.0, bold=True)
    chart_x = 0.45
    chart_y = 3.72
    bar_colors = [tokens["main"], tokens["sub"], tokens["accent"], tokens["highlight"]]
    for i, color in enumerate(bar_colors):
        bar_h = [0.45, 0.75, 0.6, 0.95][i]
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(chart_x + i * 0.42),
            Inches(chart_y + 1.0 - bar_h),
            Inches(0.26),
            Inches(bar_h),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = rgb(color)
        bar.line.color.rgb = rgb(color)

    add_pattern_box(slide, "Legend", 2.35, 3.72, 0.8, 0.28, tokens["surface_light"], tokens["border_on_light"], tokens["text_on_light"], 6)
    add_pattern_box(slide, "Note", 2.35, 4.08, 0.8, 0.28, tokens["background_light"], tokens["border_on_light"], tokens["muted_text_on_light"], 6)

    add_pattern_label(slide, "Arrow / Flow Icons", 0.45, 5.02, 2.8, bold=True)
    add_arrow_icon(slide, "Next", MSO_SHAPE.RIGHT_ARROW, 0.45, 5.35, 0.78, 0.34, tokens["main"], tokens["main"], tokens["text_on_dark"], 6)
    add_arrow_icon(slide, "Step", MSO_SHAPE.CHEVRON, 1.35, 5.35, 0.78, 0.34, tokens["sub"], tokens["sub"], tokens["text_on_dark"], 6)
    add_arrow_icon(slide, "2-way", MSO_SHAPE.LEFT_RIGHT_ARROW, 2.25, 5.35, 0.92, 0.34, tokens["accent"], tokens["accent"], tokens["text_on_dark"], 6)
    add_arrow_icon(slide, "Down", MSO_SHAPE.DOWN_ARROW, 0.45, 5.86, 0.78, 0.34, tokens["warning"], tokens["warning"], tokens["text_on_dark"], 6)
    add_arrow_icon(slide, "Up", MSO_SHAPE.UP_ARROW, 1.35, 5.86, 0.78, 0.34, tokens["success"], tokens["success"], tokens["text_on_dark"], 6)
    add_arrow_icon(slide, "Flow", MSO_SHAPE.RIGHT_ARROW, 2.25, 5.86, 0.92, 0.34, tokens["background_light"], tokens["accent"], tokens["accent"], 6)

    add_process_flow_sample(slide, 9.05, 5.85, tokens)

    add_pattern_label(slide, "Table / Grid Patterns", 3.45, 3.35, 5.8, bold=True)
    add_reusable_pattern_table(slide, "Simple", 3.45, 3.75, 1.55, 0.9, tokens, "simple")
    add_reusable_pattern_table(slide, "Header", 5.25, 3.75, 1.55, 0.9, tokens, "header")
    add_reusable_pattern_table(slide, "Horizontal", 7.05, 3.75, 1.55, 0.9, tokens, "horizontal")
    add_reusable_pattern_table(slide, "Zebra", 3.45, 5.15, 1.55, 0.9, tokens, "zebra")
    add_reusable_pattern_table(slide, "Highlight", 5.25, 5.15, 1.55, 0.9, tokens, "highlight")
    add_reusable_pattern_table(slide, "Minimal", 7.05, 5.15, 1.55, 0.9, tokens, "minimal")

    add_pattern_label(slide, "Theme Tokens", 9.65, 3.55, 2.8, bold=True)
    chips = [
        ("Main", tokens["main"]),
        ("Sub", tokens["sub"]),
        ("Accent", tokens["accent"]),
        ("Surface", tokens["surface_light"]),
        ("Dark", tokens["background_dark"]),
        ("Border", tokens["border_on_light"]),
        ("Warn", tokens["warning"]),
        ("OK", tokens["success"]),
    ]
    for i, (name, color) in enumerate(chips):
        cx = 9.65 + (i % 2) * 1.35
        cy = 3.95 + (i // 2) * 0.42
        add_pattern_box(
            slide,
            name,
            cx,
            cy,
            1.1,
            0.28,
            color,
            color,
            tokens["text_on_dark"] if color != tokens["surface_light"] else tokens["text_on_light"],
            6,
            True,
        )

    add_font_style_samples(slide, 3.15, 6.55, tokens, theme_styles)

    return slide

def build_slide(prs, slide_spec, theme_styles):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    slide_id = slide_spec["slide_id"]
    treatment = slide_spec.get("background_treatment", {})
    bg_mode = treatment.get("mode")
    bg_opacity = treatment.get("background_opacity", 1.0)

    elements = resolve_slide_elements(slide_spec)

    # 1. reference_bg を最下層に配置
    for el in elements:
        if el.get("kind") == "image" and el.get("name") == "reference_bg":
            if bg_mode == "fade_image":
                faded_path = create_faded_background(BASE_DIR / el["path"], slide_id, bg_opacity)
                add_image(slide, el, img_path_override=faded_path)
            else:
                add_image(slide, el)

    # 2. use_in_pptx=true のAssetsを背景の上に配置
    add_pptx_assets(slide, slide_spec)

    # 3. その他要素を順番に配置
    for el in elements:
        kind = el.get("kind")
        name = el.get("name")

        if kind == "image" and name != "reference_bg":
            add_image(slide, el)

        elif kind == "shape":
            add_shape(slide, el)

        elif kind == "text":
            add_textbox(slide, el, theme_styles)

    return slide

def main():
    manifest_path = BASE_DIR / "json" / "deck_manifest.json"
    defaults_path = BASE_DIR / "json" / "sector_defaults.json"

    manifest = load_json(manifest_path)

    if defaults_path.exists():
        defaults = load_json(defaults_path)
    else:
        defaults = {"theme": {"styles": {}}}

    theme_styles = defaults.get("theme", {}).get("styles", {})

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_info in manifest["slides"]:
        spec_path = BASE_DIR / slide_info["rebuild_spec"]
        if not spec_path.exists():
            print(f"[WARN] spec not found: {spec_path}")
            continue

        slide_spec = load_json(spec_path)
        build_slide(prs, slide_spec, theme_styles)
        print(f"built: {slide_spec['slide_id']}")

    add_theme_pattern_sheet(prs, theme_styles)
    print("built: theme_pattern_sheet")

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_path = BASE_DIR / "output" / f"ai_deck_reconstructor_poc_{timestamp}.pptx"
    prs.save(out_path)

    print(f"saved: {out_path}")

if __name__ == "__main__":
    main()
